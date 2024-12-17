import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import tempfile

def pdf_to_pptx(pdf_path, pptx_output_path):
    """
    Converts a PDF file into a PowerPoint presentation (PPTX).
    Each page of the PDF will be added as an image to the PPTX.

    Args:
        pdf_path (str): Path to the input PDF file.
        pptx_output_path (str): Path to save the output PPTX file.
    """
    # Convert PDF to a list of images
    print("Converting PDF pages to images...")
    try:
        images = convert_from_path(pdf_path)
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        return

    # Initialize a PowerPoint Presentation
    print("Creating PowerPoint presentation...")
    presentation = Presentation()

    # Use temporary directory for storing images
    with tempfile.TemporaryDirectory() as temp_dir:
        # Add each image as a slide
        for idx, image in enumerate(images):
            try:
                print(f"Adding page {idx + 1} to presentation...")
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                image_path = os.path.join(temp_dir, f"temp_image_{idx}.png")
                image.save(image_path, "PNG")

                # Add image to slide
                left = Inches(0)
                top = Inches(0)
                height = Inches(7.5)  # Default height for the slide
                slide.shapes.add_picture(image_path, left, top, height=height)
            except Exception as e:
                print(f"Error adding image to slide: {e}")

    # Save the PowerPoint presentation
    try:
        print("Saving the presentation...")
        presentation.save(pptx_output_path)
        print(f"Presentation saved as {pptx_output_path}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

# Example Usage
if __name__ == "__main__":
    pdf_path = "input.pdf"  # Replace with your PDF file path
    pptx_output_path = "output.pptx"  # Replace with your desired output PPTX file path
    pdf_to_pptx(pdf_path, pptx_output_path)
